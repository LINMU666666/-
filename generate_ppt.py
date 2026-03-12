"""
完整复刻：Schematic Mechanism for Synchronous NOₓ Removal and CO Oxidation over Mn-Fe/N-AC Catalyst
使用 python-pptx 生成，全程 Unicode 化学符号，无乱码
运行：python3 generate_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ══════════════════════════════════════════════════════════════════════════════
#  颜色常量（对照原图取色）
# ══════════════════════════════════════════════════════════════════════════════
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
NEAR_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
BORDER_BLUE  = RGBColor(0x2E, 0x75, 0xB6)   # 全局边框蓝
ARROW_BLUE   = RGBColor(0x1F, 0x60, 0xA8)   # 箭头蓝

# 四个面板背景
BG_TL = RGBColor(0xBD, 0xDF, 0xF7)          # 左上：浅蓝
BG_TR = RGBColor(0xFC, 0xE9, 0xC4)          # 右上：浅橙
BG_BL = RGBColor(0xBD, 0xDF, 0xF7)          # 左下：浅蓝
BG_BR = RGBColor(0xF8, 0xD4, 0xE0)          # 右下：浅粉

# 粒子 / 分子颜色
C_MN    = RGBColor(0x8B, 0x45, 0x8B)   # Mn 紫
C_FE    = RGBColor(0xCD, 0x73, 0x2A)   # Fe 橙棕
C_MN3O4 = RGBColor(0x9B, 0x55, 0x9B)   # Mn₃O₄ 深紫
C_FE3O4 = RGBColor(0xD0, 0x75, 0x30)   # Fe₃O₄ 深橙
C_N     = RGBColor(0x36, 0x8F, 0xC8)   # N原子 蓝
C_NH3   = RGBColor(0x47, 0x90, 0xCC)   # NH₃ 蓝
C_NH4   = RGBColor(0x55, 0x88, 0xCC)   # NH₄⁺ 蓝灰
C_H4P   = RGBColor(0x60, 0x90, 0xD0)   # H₄⁺ 蓝
C_NO    = RGBColor(0xA0, 0x60, 0xB0)   # NO 紫
C_CO    = RGBColor(0xBB, 0x1C, 0x1C)   # CO 深红
C_CO2   = RGBColor(0x99, 0x15, 0x15)   # CO₂ 暗红
C_O2    = RGBColor(0x2B, 0x99, 0x40)   # O₂ 绿
C_OADS  = RGBColor(0xCC, 0x55, 0x10)   # O_ads 橙红
C_FEION = RGBColor(0xB8, 0x6A, 0x28)   # Fe³⁺ 橙褐
C_BACID = RGBColor(0xC8, 0x9A, 0x30)   # B-acid 金黄
C_NAC   = RGBColor(0x86, 0x86, 0x86)   # N-AC 平台灰
C_HONEY = RGBColor(0x7A, 0x7A, 0x7A)   # 蜂巢骨架灰
C_INNER = RGBColor(0xB8, 0xB8, 0xB8)   # 蜂巢内腔浅灰
C_VAC   = RGBColor(0xFA, 0xFA, 0xE0)   # Oxygen vacancy 框背景

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


# ══════════════════════════════════════════════════════════════════════════════
#  辅助工具函数
# ══════════════════════════════════════════════════════════════════════════════

def _set_rounded(shape, radius=20000):
    """给矩形设置圆角 (radius: EMU, 建议 20000~100000)"""
    sp = shape._element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is None:
        return
    prstGeom.set('prst', 'roundRect')
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    # 清空旧值
    for gd in avLst.findall(qn('a:gd')):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {radius}')


def add_rect(slide, l, t, w, h,
             fill=None, line_rgb=None, lw=1.0, rounded=False):
    """添加矩形（可选圆角和边框）"""
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    if rounded:
        _set_rounded(shp, 35000)
    return shp


def add_oval(slide, cx, cy, r, fill, line_rgb=None, lw=0.5):
    """以中心坐标和半径添加圆形（单位：英寸）"""
    shp = slide.shapes.add_shape(
        9, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp


def add_label(slide, cx, cy, r, text, font_size=7, color=WHITE,
              bold=True, font_name='Arial'):
    """在圆形上添加居中文字标签"""
    tb = slide.shapes.add_textbox(
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_particle(slide, cx, cy, r, text, fill,
                 font_size=7, font_color=WHITE, bold=True):
    """组合：画圆 + 居中标签"""
    add_oval(slide, cx, cy, r, fill)
    add_label(slide, cx, cy, r, text, font_size=font_size,
              color=font_color, bold=bold)


def add_text(slide, l, t, w, h, text,
             font_size=9, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT,
             font_name='Arial', wrap=True):
    """单段文本框"""
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_multiline(slide, l, t, w, h, lines,
                  font_size=9, bold=False, color=BLACK,
                  align=PP_ALIGN.LEFT, font_name='Arial'):
    """多行文本框（列表行）"""
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_platform(slide, l, t, w, h=0.38, text='N-AC',
                 font_size=9):
    """灰色 N-AC 催化剂平台"""
    shp = add_rect(slide, l, t, w, h, fill=C_NAC)
    _set_rounded(shp, 15000)
    add_text(slide, l, t + 0.04, w, h - 0.04, text,
             font_size=font_size, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)


def add_arrow_line(slide, x1, y1, x2, y2,
                   color=ARROW_BLUE, lw=1.5):
    """直线连接器（带箭头）"""
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(lw)
    # 设置终点箭头
    ln = conn.line._ln
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = ln.find(qn('a:headEnd'))
    if headEnd is None:
        headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'triangle')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    return conn


# ══════════════════════════════════════════════════════════════════════════════
#  各区域绘制函数
# ══════════════════════════════════════════════════════════════════════════════

# 布局常量（英寸）
TL_L, TL_T, TL_W, TL_H = 0.10, 0.78, 4.32, 3.28   # 左上面板
TR_L, TR_T, TR_W, TR_H = 8.91, 0.78, 4.32, 3.28   # 右上面板
BL_L, BL_T, BL_W, BL_H = 0.10, 4.16, 4.32, 3.24   # 左下面板
BR_L, BR_T, BR_W, BR_H = 8.91, 4.16, 4.32, 3.24   # 右下面板
CT_L, CT_T, CT_W, CT_H = 4.55, 0.78, 4.25, 6.62   # 中心区域


def draw_title(slide):
    """标题栏（白底黑字，与原图一致）"""
    add_rect(slide, 0.05, 0.05, 13.23, 0.65,
             fill=NEAR_WHITE, line_rgb=BLACK, lw=1.5)
    add_text(slide, 0.15, 0.08, 13.03, 0.60,
             "Schematic Mechanism for Synchronous NO\u2093 Removal "
             "and CO Oxidation over Mn-Fe/N-AC Catalyst",
             font_size=17, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER, wrap=False)


def draw_outer_border(slide):
    """全幻灯片外框（蓝色粗边）"""
    add_rect(slide, 0.05, 0.05, 13.23, 7.40,
             fill=None, line_rgb=BORDER_BLUE, lw=2.5)


def draw_panel_tl(slide):
    """左上面板：NH₃ Adsorption and Activation"""
    # 面板背景（圆角）
    bg = add_rect(slide, TL_L, TL_T, TL_W, TL_H,
                  fill=BG_TL, line_rgb=BORDER_BLUE, lw=1.5, rounded=True)

    # 面板标题
    add_text(slide, TL_L + 0.05, TL_T + 0.05, TL_W - 0.1, 0.30,
             "NH\u2083 Adsorption and Activation",
             font_size=11, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # ── NH₃ 分子（3个蓝色圆，左侧竖排）
    nh3_positions = [(0.55, 1.22), (0.55, 1.68), (0.62, 2.14)]
    for cx, cy in nh3_positions:
        add_particle(slide, cx, cy, 0.21, "NH\u2083", C_NH3, font_size=7)

    # L-NH₃ 标注
    add_text(slide, 0.82, 1.55, 0.82, 0.24, "L-NH\u2083",
             font_size=8.5, bold=False, color=BLACK)

    # 箭头：NH₃ → Mn/Fe
    add_arrow_line(slide, 0.78, 1.85, 1.22, 2.30,
                   color=BORDER_BLUE, lw=1.2)

    # Mn/Fe 复合粒子（紫色大球 + 橙色中球）
    add_particle(slide, 1.45, 2.43, 0.30, "Mn/Fe", C_MN, font_size=7)
    add_particle(slide, 1.95, 2.43, 0.24, "Fe", C_FE, font_size=7)
    # L-acid 标签
    add_text(slide, 1.18, 2.76, 0.90, 0.22, "L-acid",
             font_size=8, color=BLACK, align=PP_ALIGN.CENTER)

    # Oxygen vacancy 框（带边框矩形，内含 O₂）
    add_rect(slide, 1.80, 1.08, 1.52, 0.65,
             fill=C_VAC, line_rgb=BLACK, lw=1.2)
    add_text(slide, 1.84, 1.10, 1.44, 0.26, "Oxygen vacancy",
             font_size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_particle(slide, 2.58, 1.51, 0.20, "O\u2082", C_O2, font_size=7)

    # H₄⁺ 和 NH₄⁺ 蓝色粒子（中间区域）
    add_particle(slide, 2.40, 2.28, 0.22, "H\u2084\u207A", C_H4P, font_size=7)
    add_particle(slide, 3.02, 2.08, 0.24, "NH\u2084\u207A", C_NH4, font_size=7)

    # B-acid 金黄粒子（平台右侧）
    add_particle(slide, 3.70, 2.42, 0.28, "B-acid", C_BACID, font_size=7)

    # 箭头：Oxygen vacancy → Mn/Fe
    add_arrow_line(slide, 2.30, 1.40, 1.78, 2.15,
                   color=BORDER_BLUE, lw=1.2)

    # N-AC 平台
    add_platform(slide, TL_L + 0.08, TL_T + 2.60, TL_W - 0.18)

    # 底部说明
    add_text(slide, TL_L + 0.05, TL_T + 3.04, TL_W - 0.1, 0.22,
             "Oxygen vacancies but in all oxygen vacantly",
             font_size=7.5, italic=True, color=DARK_GRAY,
             align=PP_ALIGN.LEFT)


def draw_panel_tr(slide):
    """右上面板：Mn-Fe Synergistic Redox Cycle"""
    # 面板背景
    add_rect(slide, TR_L, TR_T, TR_W, TR_H,
             fill=BG_TR, line_rgb=C_FEION, lw=1.5, rounded=True)

    # 面板标题
    add_text(slide, TR_L + 0.05, TR_T + 0.05, TR_W - 0.1, 0.28,
             "Mn-Fe Synergistic Redox Cycle",
             font_size=11, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # 反应方程核心区域（中心方框）
    add_rect(slide, TR_L + 0.10, TR_T + 0.40, TR_W - 0.20, 1.70,
             fill=RGBColor(0xFF, 0xF8, 0xEA),
             line_rgb=RGBColor(0xD0, 0x80, 0x20), lw=1.0, rounded=True)

    # 方程式
    add_text(slide, TR_L + 0.15, TR_T + 0.52, TR_W - 0.30, 0.36,
             "Mn\u2074\u207A + Fe\u00B2\u207A",
             font_size=13, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    add_text(slide, TR_L + 2.35, TR_T + 0.52, 1.80, 0.36,
             "Mn\u00B3\u207A + Fe\u00B3\u207A",
             font_size=13, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # 双向等号（⇌）
    add_text(slide, TR_L + 1.68, TR_T + 0.50, 0.62, 0.40,
             "\u21CC",
             font_size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # 上弧箭头（蓝色圆弧，模拟循环上半部 → e⁻）
    # 用半椭圆线段近似：顶部椭圆（只留上半边框）+ 文字
    # 上方循环箭头（上行 e⁻）：用大圆弧轮廓
    # 这里用一个空心圆来近似循环符号
    circ_cx = TR_L + TR_W / 2
    circ_cy = TR_T + 1.08
    circ_r  = 0.72
    shp = slide.shapes.add_shape(
        9,  # oval
        Inches(circ_cx - circ_r), Inches(circ_cy - circ_r * 0.65),
        Inches(2 * circ_r), Inches(2 * circ_r * 0.65)
    )
    shp.fill.background()
    shp.line.color.rgb = ARROW_BLUE
    shp.line.width = Pt(2.0)

    # e⁻ 标签（上方 & 下方）
    add_text(slide, TR_L + 0.30, TR_T + 0.44, 0.50, 0.28,
             "e\u207B", font_size=11, bold=True, color=DARK_GRAY)
    add_text(slide, TR_L + 3.45, TR_T + 0.44, 0.50, 0.28,
             "e\u207B", font_size=11, bold=True, color=DARK_GRAY)
    add_text(slide, TR_L + 3.25, TR_T + 1.28, 0.90, 0.28,
             "O\u2090\u1D48\u209B", font_size=10, bold=True,
             color=C_OADS)

    # 弯箭头标注（用连接线替代）
    add_arrow_line(slide, TR_L + 0.62, TR_T + 0.68,
                   TR_L + 1.62, TR_T + 0.55,
                   color=ARROW_BLUE, lw=1.5)
    add_arrow_line(slide, TR_L + 3.72, TR_T + 1.28,
                   TR_L + 2.72, TR_T + 1.44,
                   color=ARROW_BLUE, lw=1.5)

    # 说明文字
    add_multiline(slide, TR_L + 0.05, TR_T + 2.20, TR_W - 0.10, 0.95,
                  ["Electron transfer forms highly active Mn\u2074\u207A",
                   "oxidation sites and simultaneously activates",
                   "surface adsorbed oxygen species (O\u2090\u1D48\u209B)"],
                  font_size=8.5, color=DARK_GRAY)


def draw_panel_bl(slide):
    """左下面板：Low-Temperature SCR Reaction Path (L-H & E-R Mechanisms)"""
    # 面板背景
    add_rect(slide, BL_L, BL_T, BL_W, BL_H,
             fill=BG_BL, line_rgb=BORDER_BLUE, lw=1.5, rounded=True)

    # 面板标题（两行）
    add_text(slide, BL_L + 0.05, BL_T + 0.04, BL_W - 0.1, 0.28,
             "Low-Temperature SCR Reaction Path",
             font_size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, BL_L + 0.05, BL_T + 0.30, BL_W - 0.1, 0.26,
             "(L-H & E-R Mechanisms)",
             font_size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # 中间竖线（分隔 E-R 和 L-H）
    divx = BL_L + BL_W / 2
    add_rect(slide, divx - 0.01, BL_T + 0.60, 0.02, 2.34,
             fill=BLACK, line_rgb=None)

    # E-R 标题
    add_text(slide, BL_L + 0.08, BL_T + 0.62, 2.05, 0.26,
             "E-R mechanism",
             font_size=9, bold=True, color=BLACK)
    # L-H 标题
    add_text(slide, divx + 0.05, BL_T + 0.62, 2.05, 0.26,
             "L-H mechanism",
             font_size=9, bold=True, color=BLACK)

    # ── E-R 侧（左半）
    # NO 粒子（×2）
    add_particle(slide, BL_L + 0.42, BL_T + 1.02, 0.23, "NO", C_NO, font_size=8)
    add_particle(slide, BL_L + 0.42, BL_T + 1.66, 0.23, "NO", C_NO, font_size=8)

    # L-NH₃ 标签（两处）
    add_text(slide, BL_L + 0.70, BL_T + 0.96, 0.85, 0.24,
             "L-NH\u2083", font_size=8.5, color=BLACK)
    add_text(slide, BL_L + 0.08, BL_T + 1.46, 0.85, 0.24,
             "L-NH\u2083", font_size=8.5, color=BLACK)

    # 箭头（NO → Mn/Fe）
    add_arrow_line(slide, BL_L + 0.42, BL_T + 1.28, BL_L + 0.80, BL_T + 1.58,
                   color=BORDER_BLUE, lw=1.2)
    add_arrow_line(slide, BL_L + 0.42, BL_T + 1.88, BL_L + 0.80, BL_T + 2.10,
                   color=BORDER_BLUE, lw=1.2)

    # Mn/Fe 和 Mn 催化剂粒子
    add_particle(slide, BL_L + 1.05, BL_T + 2.18, 0.28, "Mn/Fe", C_MN, font_size=6.5)
    add_particle(slide, BL_L + 1.62, BL_T + 2.18, 0.24, "Mn", C_MN, font_size=8)

    # ── L-H 侧（右半）
    # Co-adsorbed L-NH₃ 标注
    add_text(slide, divx + 0.05, BL_T + 0.92, 2.00, 0.42,
             "Co-adsorbed\nL-NH\u2083",
             font_size=8, color=BLACK)
    # NO 粒子（红橙色）
    add_particle(slide, divx + 2.00, BL_T + 1.00, 0.24, "NO",
                 RGBColor(0xDD, 0x44, 0x22), font_size=7)
    # N₂ + H₂O 产物
    add_text(slide, divx + 1.75, BL_T + 1.28, 0.80, 0.45,
             "N\u2082\n+H\u2082O",
             font_size=9, bold=True,
             color=RGBColor(0x00, 0x70, 0x00), align=PP_ALIGN.CENTER)
    # NOₓ 中间体 + 结构式
    add_text(slide, divx + 0.55, BL_T + 1.72, 1.50, 0.30,
             "NO\u2093", font_size=8.5, color=DARK_GRAY)
    add_text(slide, divx + 0.55, BL_T + 2.00, 1.50, 0.28,
             "O\u2550N\u2500O",
             font_size=9, color=DARK_GRAY)

    # N-AC 平台
    add_platform(slide, BL_L + 0.08, BL_T + 2.62, BL_W - 0.18)

    # 底部
    add_text(slide, BL_L + 0.05, BL_T + 3.03, BL_W - 0.10, 0.20,
             "L-H (co-adsorption)",
             font_size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


def draw_panel_br(slide):
    """右下面板：CO/HCHO Oxidation Path (Mars-van Krevelen Mechanism)"""
    # 面板背景
    add_rect(slide, BR_L, BR_T, BR_W, BR_H,
             fill=BG_BR, line_rgb=RGBColor(0xC0, 0x40, 0x60), lw=1.5, rounded=True)

    # 面板标题（两行）
    add_text(slide, BR_L + 0.05, BR_T + 0.04, BR_W - 0.10, 0.28,
             "CO/HCHO Oxidation Path",
             font_size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, BR_L + 0.05, BR_T + 0.30, BR_W - 0.10, 0.26,
             "(Mars-van Krevelen Mechanism)",
             font_size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # CO 分子（左侧，×4）
    co_pts = [(9.25, 4.85), (9.55, 5.15), (9.25, 5.45), (9.55, 5.75)]
    for cx, cy in co_pts:
        add_particle(slide, cx, cy, 0.24, "CO", C_CO, font_size=8)

    # 弯曲箭头（CO → CO₂），用直线近似
    add_arrow_line(slide, 9.80, 4.85, 11.00, 4.62, color=ARROW_BLUE, lw=1.5)
    add_arrow_line(slide, 9.80, 5.15, 11.00, 4.85, color=ARROW_BLUE, lw=1.5)
    add_arrow_line(slide, 9.80, 5.75, 12.20, 5.75, color=ARROW_BLUE, lw=1.5)

    # CO₂ 分子（右上，×3）
    co2_pts = [(11.28, 4.62), (12.00, 4.62), (12.65, 5.12)]
    for cx, cy in co2_pts:
        add_particle(slide, cx, cy, 0.26, "CO\u2082", C_CO2, font_size=7)

    # O_ads 粒子
    add_particle(slide, 10.68, 5.32, 0.24, "O\u2090\u1D48\u209B", C_OADS,
                 font_size=6.5)

    # 右侧：O₂ 和 Fe³⁺
    add_particle(slide, 12.65, 5.72, 0.24, "O\u2082", C_O2, font_size=8)
    add_particle(slide, 13.10, 6.18, 0.24, "Fe\u00B3\u207A", C_FEION, font_size=7)

    # 循环箭头（O₂ → O_ads 回路）
    add_arrow_line(slide, 12.42, 5.72, 10.90, 5.52, color=ARROW_BLUE, lw=1.3)
    add_arrow_line(slide, 12.88, 5.98, 12.65, 5.99, color=ARROW_BLUE, lw=1.3)

    # Mn/Fe 粒子（×2，在 N-AC 平台上）
    add_particle(slide, 9.62, 6.22, 0.28, "Mn/Fe", C_MN, font_size=6.5)
    add_particle(slide, 10.22, 6.22, 0.28, "Mn/Fe", C_MN, font_size=6.5)

    # N-AC 平台
    add_platform(slide, BR_L + 0.08, BR_T + 2.60, BR_W - 0.18)


def draw_center(slide):
    """中心：Mn-Fe/N-AC 蜂窝骨架结构 + 粒子"""
    # 中心区域背景（浅灰）
    bg = add_rect(slide, CT_L, CT_T, CT_W, CT_H,
                  fill=RGBColor(0xE0, 0xE0, 0xE0),
                  line_rgb=RGBColor(0xAA, 0xAA, 0xAA), lw=1.0)
    _set_rounded(bg, 20000)

    # ── 蜂窝/六边形网格（用矩形阵列模拟横截面孔道）
    cell_w, cell_h = 0.58, 0.56
    grid_l = CT_L + 0.10
    grid_t = CT_T + 0.28
    cols, rows = 6, 9
    for r in range(rows):
        for c in range(cols):
            cl = grid_l + c * cell_w
            ct = grid_t + r * cell_h
            # 骨架（深灰）
            cell = slide.shapes.add_shape(
                1, Inches(cl), Inches(ct),
                Inches(cell_w - 0.04), Inches(cell_h - 0.04)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_HONEY
            cell.line.fill.background()
            # 孔道内腔（浅灰）
            im = 0.07
            inner = slide.shapes.add_shape(
                1,
                Inches(cl + im), Inches(ct + im),
                Inches(cell_w - 0.04 - 2 * im),
                Inches(cell_h - 0.04 - 2 * im)
            )
            inner.fill.solid()
            inner.fill.fore_color.rgb = C_INNER
            inner.line.fill.background()

    # ── 粒子（Mn₃O₄、Fe₃O₄、N 原子）
    # Mn₃O₄ 紫色
    mn3o4_pts = [
        (5.40, 1.22, 0.30),  # top area
        (5.28, 2.62, 0.26),
        (7.20, 4.82, 0.22),
    ]
    for cx, cy, r in mn3o4_pts:
        add_particle(slide, cx, cy, r,
                     "Mn\u2083O\u2084", C_MN3O4, font_size=6)

    # Fe₃O₄ 橙色
    fe3o4_pts = [
        (7.10, 1.28, 0.30),  # top area
        (6.85, 3.05, 0.26),
        (5.95, 4.35, 0.22),
        (6.90, 5.85, 0.24),
    ]
    for cx, cy, r in fe3o4_pts:
        add_particle(slide, cx, cy, r,
                     "Fe\u2083O\u2084", C_FE3O4, font_size=6)

    # N 原子（蓝色小圆，散布在骨架上）
    n_pts = [
        (5.82, 1.08), (6.48, 1.14), (5.62, 1.60), (6.88, 1.52),
        (5.20, 2.02), (6.20, 1.98), (7.42, 1.90), (5.55, 2.45),
        (7.05, 2.35), (5.00, 3.20), (6.40, 2.88), (7.50, 3.05),
        (5.30, 3.75), (6.70, 3.62), (5.80, 4.18), (7.30, 4.22),
        (5.10, 5.00), (6.45, 5.28), (7.68, 5.42), (5.55, 5.85),
        (7.05, 6.35), (6.20, 6.70),
    ]
    for nx, ny in n_pts:
        add_particle(slide, nx, ny, 0.14, "N", C_N, font_size=6.5)

    # ── 顶部标签 Mn₃O₄ / Fe₃O₄
    add_text(slide, CT_L + 0.05, CT_T + 0.03, 1.90, 0.28,
             "Mn\u2083O\u2084",
             font_size=13, bold=True, color=C_MN3O4,
             align=PP_ALIGN.LEFT)
    add_text(slide, CT_L + 2.20, CT_T + 0.03, 1.90, 0.28,
             "Fe\u2083O\u2084",
             font_size=13, bold=True, color=C_FE3O4,
             align=PP_ALIGN.LEFT)

    # 底部 Mn-Fe/N-AC 标签
    add_text(slide, CT_L + 0.10, CT_T + CT_H - 0.42, CT_W - 0.20, 0.38,
             "Mn-Fe/N-AC",
             font_size=14, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)


def draw_annotations(slide):
    """面板 ↔ 中心 的标注文字（带箭头）"""

    # ── 1. N-doping → 左上角（从中心上方引出）
    add_text(slide, CT_L - 0.05, CT_T + 0.42, 4.00, 0.55,
             "N-doping enhances surface\nbasicity for NH\u2083 capture.",
             font_size=9, italic=True, color=DARK_GRAY,
             align=PP_ALIGN.LEFT)
    add_arrow_line(slide, CT_L + 0.10, CT_T + 0.68,
                   TL_L + TL_W - 0.05, TL_T + 0.50,
                   color=BORDER_BLUE, lw=1.3)

    # ── 2. Synergy enables SCR（中心 → 左下）
    add_text(slide, CT_L - 0.05, CT_T + CT_H - 2.10, 4.10, 0.55,
             "Synergy enables\nefficient SCR at 280\u00B0C.",
             font_size=9, italic=True, color=DARK_GRAY,
             align=PP_ALIGN.LEFT)
    add_arrow_line(slide, CT_L + 0.10, CT_T + CT_H - 1.65,
                   BL_L + BL_W - 0.05, BL_T + 0.60,
                   color=BORDER_BLUE, lw=1.3)

    # ── 3. Mn-Fe cycle accelerates → 右下
    add_text(slide, CT_L + CT_W - 0.05, CT_T + CT_H - 2.10, 4.10, 0.70,
             "Mn-Fe cycle accelerates\noxygen activation for\nCO oxidation",
             font_size=9, italic=True, color=DARK_GRAY,
             align=PP_ALIGN.LEFT)
    add_arrow_line(slide, CT_L + CT_W - 0.05, CT_T + CT_H - 1.65,
                   BR_L + 0.10, BR_T + 0.60,
                   color=BORDER_BLUE, lw=1.3)

    # ── 4. 面板 → 中心的连接线（4条）
    # 左上面板 → 中心
    add_arrow_line(slide, TL_L + TL_W, TL_T + TL_H / 2,
                   CT_L, CT_T + CT_H / 4,
                   color=ARROW_BLUE, lw=1.8)
    # 右上面板 → 中心
    add_arrow_line(slide, TR_L, TR_T + TR_H / 2,
                   CT_L + CT_W, CT_T + CT_H / 4,
                   color=ARROW_BLUE, lw=1.8)
    # 左下面板 → 中心
    add_arrow_line(slide, BL_L + BL_W, BL_T + BL_H / 2,
                   CT_L, CT_T + CT_H * 3 / 4,
                   color=ARROW_BLUE, lw=1.8)
    # 右下面板 → 中心
    add_arrow_line(slide, BR_L, BR_T + BR_H / 2,
                   CT_L + CT_W, CT_T + CT_H * 3 / 4,
                   color=ARROW_BLUE, lw=1.8)


# ══════════════════════════════════════════════════════════════════════════════
#  主函数
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    # 白色背景
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # 绘制各部分（顺序很重要：先背景，后前景）
    draw_outer_border(slide)
    draw_title(slide)
    draw_center(slide)       # 中心先画（做底层）
    draw_panel_tl(slide)
    draw_panel_tr(slide)
    draw_panel_bl(slide)
    draw_panel_br(slide)
    draw_annotations(slide)  # 最后画标注

    return prs


def main():
    prs = build_presentation()
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Mn-Fe_N-AC_Catalyst_Mechanism.pptx")
    prs.save(out)
    print(f"[OK] 已生成 => {out}")
    print(f"     幻灯片数: {len(prs.slides)}")
    print(f"     形状数:   {len(prs.slides[0].shapes)}")


if __name__ == "__main__":
    main()
